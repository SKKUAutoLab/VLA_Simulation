#!/usr/bin/env python3
"""
restyle_arch_slide.py — 사용자가 5페이지에 그린 '박스 다이어그램' 스타일로
VLA 아키텍처 흐름도 슬라이드를 제자리(in-place) 교체한다.

사용자 스타일(5페이지에서 추출):
  - 노드 박스: 둥근 사각형, 채움 없음, 외곽선 네이비 #1D3155, 맑은 고딕 10pt
  - 패키지 박스: 둥근 사각형, 채움 없음, 외곽선 검정, 제목 맑은 고딕 16pt(좌상단)
  - 제어/출력 강조: 분홍 채움 #FECAE8
  - 화살표: 검은 직선 연결선(삼각 화살촉)
  - 폰트: 맑은 고딕 (latin + ea)

주의: 파일을 그대로 열어 해당 슬라이드만 수정하므로 사용자가 추가한 5페이지 등 다른 슬라이드는 보존된다.
"""
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

NAVY  = RGBColor(0x1D, 0x31, 0x55)   # 노드 외곽선
BLACK = RGBColor(0x00, 0x00, 0x00)
PINK  = RGBColor(0xFE, 0xCA, 0xE8)   # 제어 강조
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY_TXT = RGBColor(0xE6, 0xE6, 0xE6)  # 회색 노드의 보조 글자
KFONT = "맑은 고딕"

# 역할별 패키지 색 (사용자 5페이지와 동일: 테마색 + 명암보정)
ROLE = {
    "sensor":   ("accent4", 20000, 80000),   # 입력/센서 = 노랑(Perception)
    "vla":      ("accent6", 40000, 60000),   # VLA 판단 = 초록(Decision)
    "sim":      ("accent1", 20000, 80000),   # 시뮬/제어브리지 = 파랑
}


def kfont(run, size, bold=False, color=BLACK):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set('typeface', KFONT)


def _clear_fill(spPr):
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        for e in spPr.findall(qn(tag)):
            spPr.remove(e)


def _put_fill(shape, sf):
    spPr = shape._element.spPr
    _clear_fill(spPr)
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(sf)
    else:
        spPr.append(sf)


def scheme_fill(shape, scheme, lumMod=None, lumOff=None):
    spPr = shape._element.spPr
    sf = spPr.makeelement(qn('a:solidFill'), {})
    clr = sf.makeelement(qn('a:schemeClr'), {'val': scheme})
    if lumMod is not None:
        clr.append(clr.makeelement(qn('a:lumMod'), {'val': str(lumMod)}))
    if lumOff is not None:
        clr.append(clr.makeelement(qn('a:lumOff'), {'val': str(lumOff)}))
    sf.append(clr)
    _put_fill(shape, sf)


def srgb_fill(shape, hexstr):
    spPr = shape._element.spPr
    sf = spPr.makeelement(qn('a:solidFill'), {})
    sf.append(sf.makeelement(qn('a:srgbClr'), {'val': hexstr}))
    _put_fill(shape, sf)


def _txt(shape, lines, size, bold, color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         title=False):
    tf = shape.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP if title else anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT if title else align
        r = p.add_run(); r.text = ln
        # 첫 줄(노드 파일명/제목)만 약간 크게·굵게
        if i == 0:
            kfont(r, size, bold, color)
        else:
            kfont(r, size - 1.5, False, RGBColor(0x40, 0x40, 0x40))


def pkg(s, title, x, y, w, h, role="vla"):
    """패키지 컨테이너: 둥근 사각형, 역할별 파스텔 채움, 검정 외곽선, 제목 좌상단 16pt."""
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    b.line.color.rgb = BLACK; b.line.width = Pt(1.25); b.shadow.inherit = False
    scheme_fill(b, *ROLE[role])
    _txt(b, [title], 16, False, BLACK, title=True)
    return b


def node(s, lines, x, y, w, h, pink=False, size=10):
    """노드 박스: 둥근 사각형, 회색 채움(흰 글씨) 또는 분홍(검은 글씨), 네이비 외곽선."""
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    b.line.color.rgb = NAVY; b.line.width = Pt(1.0); b.shadow.inherit = False
    if pink:
        srgb_fill(b, "FECAE8"); txt_col, sub_col = BLACK, RGBColor(0x50,0x50,0x50)
    else:
        scheme_fill(b, "bg1", lumMod=50000); txt_col, sub_col = WHITE, GREY_TXT
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        kfont(r, size if i == 0 else size-1.5, i == 0, txt_col if i == 0 else sub_col)
    return b


def ext(s, lines, x, y, w, h):
    """외부 시스템(Gazebo) 컨테이너: 센서 역할(노랑) 파스텔 박스."""
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    b.line.color.rgb = BLACK; b.line.width = Pt(1.25); b.shadow.inherit = False
    scheme_fill(b, *ROLE["sensor"])
    _txt(b, lines, 12, True, BLACK)
    return b


def conn(s, x1, y1, x2, y2, label=None, lx=None, ly=None, lcolor=RGBColor(0x30,0x30,0x30)):
    """검은 직선 화살표 연결선 + (옵션) 토픽 라벨."""
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = BLACK; c.line.width = Pt(1.5); c.shadow.inherit = False
    ln = c.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    if label:
        tb = s.shapes.add_textbox(Inches(lx if lx is not None else (x1+x2)/2 - 0.7),
                                  Inches(ly if ly is not None else (y1+y2)/2 - 0.18),
                                  Inches(1.7), Inches(0.34))
        tf = tb.text_frame; tf.word_wrap = False
        for i, t in enumerate(label.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = t; kfont(r, 8.5, False, lcolor)
    return c


def clear_nonplaceholder(slide):
    for sh in list(slide.shapes):
        if not sh.is_placeholder:
            sh._element.getparent().remove(sh._element)


def build(slide):
    s = slide
    clear_nonplaceholder(s)

    # ── 범례 (우상단) — 사용자 5페이지 형식(색 스와치 + 라벨) ──────────
    def swatch(text, role, x, pink=False):
        b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.1), Inches(0.3), Inches(0.26))
        b.line.color.rgb = BLACK; b.line.width = Pt(0.75); b.shadow.inherit = False
        if pink: srgb_fill(b, "FECAE8")
        else: scheme_fill(b, *ROLE[role])
        tb = s.shapes.add_textbox(Inches(x+0.33), Inches(1.08), Inches(1.6), Inches(0.3))
        r = tb.text_frame.paragraphs[0].add_run(); r.text = text; kfont(r, 10, False, BLACK)
    swatch("센서/입력", "sensor", 6.55)
    swatch("VLA 판단", "vla", 8.35)
    swatch("시뮬/제어", "sim", 10.1)
    swatch("제어 메시지", "", 11.9, pink=True)

    # ── 패키지/박스 배치 ──────────────────────────────────────────────
    # Gazebo (센서/입력 = 노랑, 좌측)
    ext(s, ["Gazebo 시뮬레이터",
            "camera/image_raw · scan · /odom 발행",
            "/cmd_vel 구독 → 차량 구동"], 0.3, 2.5, 3.05, 1.5)

    # lora_pipeline (VLA 판단 = 초록) — 중앙 컨테이너
    pkg(s, "lora_pipeline  (VLA 노드)", 3.65, 1.55, 6.0, 3.05, role="vla")
    node(s, ["vla_gui.py", "명령 콘솔(GUI)"],                 3.95, 2.15, 2.45, 0.72)
    node(s, ["vla_brain_node.py", "Slow · Qwen3-VL 브레인"], 3.95, 3.4, 2.45, 0.82)
    node(s, ["vla_lora_drive_node.py", "Fast · 실시간 주행"], 6.75, 2.75, 2.65, 0.9)

    # interfaces_pkg — 제어 메시지(분홍)
    node(s, ["interfaces_pkg", "MotionCommand (msg)"], 10.05, 2.8, 2.9, 0.85, pink=True)

    # simulation_pkg (시뮬/제어 = 파랑) — 하단 컨테이너
    pkg(s, "simulation_pkg", 3.65, 4.9, 6.0, 1.85, role="sim")
    node(s, ["simulation_sender_node.py", "MotionCommand → cmd_vel"], 3.95, 5.5, 2.75, 0.95)
    node(s, ["load_ego_car_node.py", "차량 스폰(ego_vehicle)"],        6.95, 5.5, 2.45, 0.95)

    # ── 화살표(토픽) ─────────────────────────────────────────────────
    conn(s, 3.35, 3.12, 6.75, 3.12, "camera/image_raw · scan · /odom", lx=4.15, ly=3.14)   # 센서→drive
    conn(s, 5.17, 2.87, 5.17, 3.4, "nl_command", lx=5.35, ly=3.0, lcolor=RGBColor(0x7A,0x5A,0x9A))  # gui→brain
    conn(s, 6.4, 3.62, 6.75, 3.3, "vla/command", lx=6.15, ly=3.78, lcolor=RGBColor(0xC0,0x5A,0x10))  # brain→drive
    conn(s, 9.4, 3.2, 10.05, 3.2, "topic_control_signal", lx=8.6, ly=2.5)   # drive→interfaces
    conn(s, 8.0, 3.65, 6.2, 5.5, "MotionCommand", lx=6.6, ly=4.4)           # drive→sender
    conn(s, 3.95, 5.75, 1.7, 4.0, "/cmd_vel", lx=1.7, ly=4.6)               # sender→Gazebo
    conn(s, 7.5, 3.65, 5.6, 4.1, "vla/cur_lane", lx=6.05, ly=4.0, lcolor=RGBColor(0x60,0x60,0x60))  # drive↔brain

    # 하단 캡션
    cap = s.shapes.add_textbox(Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.4))
    r = cap.text_frame.paragraphs[0].add_run()
    r.text = "명령(gui→brain→drive) → 제어(drive→sender→Gazebo) → 센서 피드백(Gazebo→drive)이 순환하는 닫힌 루프"
    kfont(r, 12, True, RGBColor(0xC0, 0x5A, 0x10))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("-f", "--file", default="/home/autolab/Downloads/VLA_교육자료.pptx")
    ap.add_argument("--find", default="흐름도", help="대상 슬라이드 제목 키워드")
    a = ap.parse_args()
    prs = Presentation(a.file)
    target = None
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and a.find in sh.text_frame.text:
                target = s; break
        if target: break
    if target is None:
        raise SystemExit(f"'{a.find}' 슬라이드를 찾지 못함")
    build(target)
    prs.save(a.file)
    print("restyled in-place:", a.file, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
