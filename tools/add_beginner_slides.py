#!/usr/bin/env python3
"""
add_beginner_slides.py — 발표자료를 재생성하지 않고(사용자 수정분 보존),
비전공자용 온보딩 슬라이드 2장을 '3. 환경 준비' 섹션 직후에 끼워 넣는다.

  A) 사전 준비물 & 요구사항 (OS/GPU/사전설치/확인)
  B) 0단계 — 코드 받기 · 파이썬 환경 · 설치 검증

기존 콘텐츠 슬라이드의 레이아웃/서식(맑은 고딕·❑ 개요)을 그대로 재사용한다.
"""
import argparse, copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

BLUE   = RGBColor(0x00, 0x70, 0xC0)
CODEBG = RGBColor(0x1E, 0x1E, 0x1E)
CODEFG = RGBColor(0xE6, 0xE6, 0xE6)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
MONO   = "Consolas"


def _title_body(s):
    """placeholder를 idx가 아닌 '타입'으로 구분 (이 템플릿은 TITLE/OBJECT가 둘 다 idx 0)."""
    title = body = None
    for ph in s.placeholders:
        t = ph.placeholder_format.type
        if t in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE) and title is None:
            title = ph
        elif t in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) and body is None:
            body = ph
    return title, body


def add_content(prs, ref_slide, title):
    """기존 콘텐츠 슬라이드(ref)의 제목·본문 placeholder를 복제해 동일 서식 유지.
    (현재 레이아웃이 'Blank Slide'라 add_slide만으론 제목/본문 틀이 안 생김)"""
    s = prs.slides.add_slide(ref_slide.slide_layout)   # 배경·푸터·로고 상속
    # 제목 틀만 복제 (본문은 명시적 텍스트박스로 그림 — Blank 레이아웃 상속 문제 회피)
    for ph in ref_slide.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            s.shapes._spTree.append(copy.deepcopy(ph._element))
            break
    tph, _ = _title_body(s)
    if tph is not None:
        tph.text_frame.text = title
        for para in tph.text_frame.paragraphs:
            for r in para.runs:
                r.font.bold = True; r.font.color.rgb = BLUE
    return s


# 템플릿 개요 룩(❑/○/·)을 명시적으로 재현 — 색/크기 고정
_BULLET = ["❑ ", "○ ", "· "]
_BSIZE  = [16, 14, 13]
_BCOLOR = [BLACK, RGBColor(0x40, 0x40, 0x40), RGBColor(0x55, 0x55, 0x55)]
_KFONT  = "맑은 고딕"


def body_box(s, items, left, top, width, height):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    items = [it for it in items if (it[0] if isinstance(it, tuple) else it).strip()]
    first = True
    for it in items:
        txt, lvl = (it if isinstance(it, tuple) else (it, 0))
        lvl = min(lvl, 2)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(5)
        p.level = lvl
        r = p.add_run()
        r.text = ("    " * lvl) + _BULLET[lvl] + txt
        r.font.size = Pt(_BSIZE[lvl]); r.font.bold = (lvl == 0)
        r.font.color.rgb = _BCOLOR[lvl]
        r.font.name = _KFONT
        rPr = r._r.get_or_add_rPr()
        for tag in ('a:latin', 'a:ea', 'a:cs'):
            e = rPr.find(qn(tag))
            if e is None:
                e = rPr.makeelement(qn(tag), {}); rPr.append(e)
            e.set('typeface', _KFONT)


def hdr(s, text, top, size=16):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.45))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = BLUE


def code_box(s, lines, top, height, left=0.55, width=12.2, size=12.5):
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = CODEBG
    box.line.color.rgb = RGBColor(0x40, 0x40, 0x40); box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(1)
        r = p.add_run(); r.text = ln if ln else " "
        r.font.name = MONO; r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0x8A, 0xA8, 0x8A) if ln.strip().startswith("#") else CODEFG


def note(s, text, top, color=ORANGE, size=12):
    tb = s.shapes.add_textbox(Inches(0.55), Inches(top), Inches(12.2), Inches(0.5))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("-f", "--file", default="/home/autolab/Downloads/VLA_교육자료.pptx")
    a = ap.parse_args()
    prs = Presentation(a.file)

    # '3. 환경 준비' 섹션 헤더 인덱스 찾기
    anchor_idx = None
    for i, s in enumerate(prs.slides):
        if s.shapes and s.shapes[0].has_text_frame and "환경 준비" in s.shapes[0].text_frame.text:
            anchor_idx = i; break
    if anchor_idx is None:
        raise SystemExit("'환경 준비' 섹션을 찾지 못함")

    # 참조 슬라이드 = 실제 ❑ 개요(불릿)를 쓰는 슬라이드 — 제목·본문 서식을 그대로 복제
    ref_slide = None
    for s in prs.slides:
        if s.shapes and s.shapes[0].has_text_frame and "핵심 원리" in s.shapes[0].text_frame.text:
            ref_slide = s; break
    if ref_slide is None:
        ref_slide = prs.slides[anchor_idx + 1]

    # ── 슬라이드 A: 사전 준비물 & 요구사항 ──
    sA = add_content(prs, ref_slide, "사전 준비물 & 요구사항 (시작 전 체크)")
    body_box(sA, [
        ("이 실습에 필요한 것 (준비물)", 0),
        ("OS: Ubuntu 22.04 LTS", 1),
        ("GPU: NVIDIA 그래픽카드 (VRAM 8GB 이상 권장) + 드라이버·CUDA", 1),
        ("메모리 16GB+ / 디스크 여유 20GB+ (모델·데이터 저장)", 1),
        ("인터넷 연결 (최초 1회 모델·패키지 다운로드)", 1),
        ("미리 설치돼 있어야 하는 것 (수업용 PC엔 이미 세팅됨)", 0),
        ("ROS 2 Humble  —  아래 확인이 되면 OK", 1),
        ("NVIDIA 드라이버 + CUDA  —  nvidia-smi 가 표가 뜨면 OK", 1),
    ], 0.42, 1.35, 12.4, 3.6)
    hdr(sA, "환경 확인 (터미널을 열고 그대로 입력)", 4.55)
    code_box(sA, [
        "ros2 --version          # 'ros2 ...' 버전이 나오면 ROS2 설치 OK",
        "nvidia-smi              # GPU 표가 나오면 드라이버/CUDA OK",
    ], top=5.0, height=1.0, size=12.5)
    note(sA, "둘 중 하나라도 안 나오면, 먼저 조교(또는 랩 담당)에게 환경 세팅을 요청하세요.", top=6.2)

    # ── 슬라이드 B: 0단계 — 코드 받기 · 파이썬 · 검증 ──
    sB = add_content(prs, ref_slide, "0단계 — 코드 받기 · 파이썬 환경 · 검증")
    hdr(sB, "① 코드 받기 (최초 1회)", 1.35)
    code_box(sB, [
        "git clone https://github.com/SKKUAutoLab/VLA_Simulation.git ~/ros2_autonomous_vehicle_simulation",
        "cd ~/ros2_autonomous_vehicle_simulation      # 이후 모든 명령은 이 폴더에서",
    ], top=1.8, height=0.95, size=11.5)
    hdr(sB, "② 파이썬 환경 주의 (conda 충돌 방지)", 2.95)
    code_box(sB, [
        "conda deactivate 2>/dev/null   # Anaconda(base) 켜져 있으면 끄기",
        "which python3                  # /usr/bin/python3 여야 함 (conda 경로면 위 명령 반복)",
    ], top=3.4, height=0.95, size=12)
    hdr(sB, "③ 설치 검증 (설치·빌드 후 확인)", 4.55)
    code_box(sB, [
        "ros2 --version",
        "python3 -c \"import torch; print('CUDA:', torch.cuda.is_available())\"   # True 면 GPU 준비 완료",
    ], top=5.0, height=0.95, size=12)
    note(sB, "여기까지 되면 다음 장 '설치 & 빌드'를 그대로 따라가면 됩니다.", top=6.15)

    # ── 삽입 위치로 이동: 섹션 헤더(anchor) 바로 뒤 ──
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    newA, newB = ids[-2], ids[-1]          # 방금 추가한 2장 (맨 끝)
    anchor = ids[anchor_idx]
    sldIdLst.remove(newA); sldIdLst.remove(newB)
    anchor.addnext(newA)                    # anchor → A
    newA.addnext(newB)                      # A → B

    prs.save(a.file)
    print("inserted 2 beginner slides after '3. 환경 준비' |", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    main()
