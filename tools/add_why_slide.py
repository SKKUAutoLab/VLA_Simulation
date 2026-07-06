#!/usr/bin/env python3
"""'왜 이 교육을 하는가 — 학습 목표' 슬라이드를 '1. VLA란 무엇인가' 다음에 삽입(제자리 편집).
템플릿 폰트/개요(❑/○) 유지. 기존 슬라이드(손그림 5페이지 포함) 보존."""
import argparse, copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn

BLUE=RGBColor(0x00,0x70,0xC0); BLACK=RGBColor(0x1A,0x1A,0x1A)
KO="맑은 고딕"; _B=["❑ ","○ ","· "]; _S=[16,14,13]
_C=[BLACK,RGBColor(0x40,0x40,0x40),RGBColor(0x55,0x55,0x55)]


def _kf(r,sz,bold,color):
    r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=color; r.font.name=KO
    rPr=r._r.get_or_add_rPr()
    for t in ('a:latin','a:ea','a:cs'):
        e=rPr.find(qn(t))
        if e is None: e=rPr.makeelement(qn(t),{}); rPr.append(e)
        e.set('typeface',KO)


def title(s, ref, text):
    for ph in ref.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE,PP_PLACEHOLDER.CENTER_TITLE):
            s.shapes._spTree.append(copy.deepcopy(ph._element)); break
    for ph in s.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE,PP_PLACEHOLDER.CENTER_TITLE):
            ph.text_frame.text=text
            for para in ph.text_frame.paragraphs:
                for r in para.runs: r.font.bold=True; r.font.color.rgb=BLUE
            break


def body(s, items, left, top, width, height):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(height))
    tb.text_frame.word_wrap=True; first=True
    for it in items:
        txt,lvl=(it if isinstance(it,tuple) else (it,0)); lvl=min(lvl,2)
        p=tb.text_frame.paragraphs[0] if first else tb.text_frame.add_paragraph(); first=False
        p.space_after=Pt(6); p.level=lvl
        r=p.add_run(); r.text=("    "*lvl)+_B[lvl]+txt; _kf(r,_S[lvl],lvl==0,_C[lvl])


def main():
    ap=argparse.ArgumentParser(); ap.add_argument("-f","--file",default="/home/autolab/Downloads/VLA_교육자료.pptx")
    a=ap.parse_args(); prs=Presentation(a.file)
    # 앵커: '1. VLA란 무엇인가' 섹션 divider
    anchor=None
    for i,s in enumerate(prs.slides):
        if s.shapes and s.shapes[0].has_text_frame and "VLA란 무엇인가" in s.shapes[0].text_frame.text:
            anchor=i; break
    if anchor is None: raise SystemExit("'1. VLA란 무엇인가' 슬라이드 못 찾음")
    ref=None
    for s in prs.slides:
        if s.shapes and s.shapes[0].has_text_frame and "핵심 원리" in s.shapes[0].text_frame.text:
            ref=s; break
    ref=ref or prs.slides[anchor]

    s=prs.slides.add_slide(ref.slide_layout)
    title(s, ref, "왜 이 교육을 하는가 — 학습 목표")
    body(s, [
        ("자율주행이 바뀌고 있다: 사람이 규칙을 짜던 방식 → 데이터로 배우는 방식(VLA)", 0),
        ("고전 방식은 모든 상황의 규칙을 사람이 설계 → 예외에 취약", 1),
        ("VLA는 사람의 주행을 모방 학습 → 처음 보는 상황·자연어 지시에 유연", 1),
        ("'말로 지시하면 움직이는' 로봇 — 언어와 행동이 이어진다", 1),
        ("이 교육에서 직접 해보며 얻는 것", 0),
        ("이미지 한 장이 어떻게 '주행 행동'이 되는지 스스로 만들어 본다", 1),
        ("데이터 수집 → 학습 → 추론 → 주행까지 전 과정을 손으로 체득", 1),
        ("20억 파라미터 거대 모델을 '조금만' 학습해 쓰는 실전법(freeze + 작은 Head)", 1),
        ("목표: 비전공자도 VLA의 동작 원리를 코드로 이해한다", 0),
    ], 0.5, 1.5, 12.3, 5.2)

    lst=prs.slides._sldIdLst; ids=list(lst); newid=ids[-1]; anchor_id=ids[anchor]
    lst.remove(newid); anchor_id.addnext(newid)
    prs.save(a.file)
    print("inserted '왜 이 교육을 하는가' after '1. VLA란 무엇인가' |", len(prs.slides._sldIdLst), "slides")


if __name__=="__main__":
    main()
