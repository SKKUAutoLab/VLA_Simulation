#!/usr/bin/env python3
"""'자동차의 행동 공간(Action Space)' 슬라이드를 '[7]-[8] 자연어 명령' 다음에 삽입.
템플릿 폰트/개요(❑/○) 유지, 기존 슬라이드 보존."""
import argparse, copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn

BLUE=RGBColor(0x00,0x70,0xC0); BLACK=RGBColor(0x1A,0x1A,0x1A); ACC=RGBColor(0xB0,0x40,0x2A)
KO="맑은 고딕"; _B=["❑ ","○ ","· "]; _S=[16,14,13]
_C=[BLACK,RGBColor(0x40,0x40,0x40),RGBColor(0x55,0x55,0x55)]


def _kf(r,sz,bold,color):
    r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=color; r.font.name=KO
    rPr=r._r.get_or_add_rPr()
    for t in ('a:latin','a:ea','a:cs'):
        e=rPr.find(qn(t))
        if e is None: e=rPr.makeelement(qn(t),{}); rPr.append(e)
        e.set('typeface',KO)


def title(s, ref, text):
    for ph in ref.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE,PP_PLACEHOLDER.CENTER_TITLE):
            s.shapes._spTree.append(copy.deepcopy(ph._element)); break
    for ph in s.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE,PP_PLACEHOLDER.CENTER_TITLE):
            ph.text_frame.text=text
            for para in ph.text_frame.paragraphs:
                for r in para.runs: r.font.bold=True; r.font.color.rgb=BLUE
            break


def body(s, items, left, top, width, height):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(height))
    tb.text_frame.word_wrap=True; first=True
    for it in items:
        txt,lvl=(it if isinstance(it,tuple) else (it,0)); lvl=min(lvl,2)
        p=tb.text_frame.paragraphs[0] if first else tb.text_frame.add_paragraph(); first=False
        p.space_after=Pt(5); p.level=lvl
        r=p.add_run(); r.text=("    "*lvl)+_B[lvl]+txt; _kf(r,_S[lvl],lvl==0,_C[lvl])


def main():
    ap=argparse.ArgumentParser(); ap.add_argument("-f","--file",default="/home/autolab/Downloads/VLA_교육자료.pptx")
    a=ap.parse_args(); prs=Presentation(a.file)
    anchor=None
    for i,s in enumerate(prs.slides):
        if s.shapes and s.shapes[0].has_text_frame and "자연어 명령" in s.shapes[0].text_frame.text:
            anchor=i; break
    if anchor is None: raise SystemExit("'자연어 명령' 슬라이드 못 찾음")
    ref=None
    for s in prs.slides:
        if s.shapes and s.shapes[0].has_text_frame and "핵심 원리" in s.shapes[0].text_frame.text:
            ref=s; break
    ref=ref or prs.slides[anchor]

    s=prs.slides.add_slide(ref.slide_layout)
    title(s, ref, "자동차의 행동 공간 (Action Space)")
    body(s, [
        ("행동 공간 = 자율주행차가 실제로 할 수 있는 ‘행동의 집합’. 자연어 명령은 이 중 하나로 변환된다.", 0),
        ("지원 행동", 0),
        ("주행(차선 N · K바퀴/무한) · 차선 변경 · 정지 · 일시정지/재개", 1),
        ("목적지 정지(랜드마크) · 방향 전환 · 복구(차선 스냅)", 1),
        ("후진 · 속도 조절(천천히/빨리) · 추월", 1),
        ("(자율·결정론적) 라이다 전방 정지 · 자동 회피", 1),
        ("헷갈리기 쉬운 차이 — 자동 회피 vs 추월", 0),
        ("자동 회피: 앞이 막히면 스스로 옆 차선으로 → 그대로 눌러앉음 (반사적 안전)", 1),
        ("추월: 옆으로 나가 앞차를 지나친 뒤 원래 차선으로 복귀 (능동 기동)", 1),
        ("핵심: 언어모델이 이해해도 ‘행동 공간’에 없으면 실행 불가 → 행동을 늘리는 것이 곧 능력 확장.", 0),
    ], 0.5, 1.45, 12.4, 5.3)

    lst=prs.slides._sldIdLst; ids=list(lst); newid=ids[-1]; anchor_id=ids[anchor]
    lst.remove(newid); anchor_id.addnext(newid)
    prs.save(a.file)
    print("inserted 'action space' slide after 자연어 명령 |", len(prs.slides._sldIdLst), "slides")


if __name__=="__main__":
    main()
