#!/usr/bin/env python3
"""'노트북 없는 학생 — 무료 브라우저 트랙' 슬라이드를 '사전 학습 산출물 배포' 다음에 삽입."""
import argparse, copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLUE=RGBColor(0x00,0x70,0xC0); ORANGE=RGBColor(0xED,0x7D,0x31)
BLACK=RGBColor(0x1A,0x1A,0x1A); GREY=RGBColor(0x55,0x55,0x55); GREEN=RGBColor(0x2E,0x7D,0x32)
KO="맑은 고딕"; _B=["❑ ","○ ","· "]; _S=[15,13.5,12.5]; _C=[BLACK,RGBColor(0x40,0x40,0x40),GREY]


def _kf(r,sz,bold,color):
    r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=color; r.font.name=KO
    rPr=r._r.get_or_add_rPr()
    for t in ('a:latin','a:ea','a:cs'):
        e=rPr.find(qn(t))
        if e is None: e=rPr.makeelement(qn(t),{}); rPr.append(e)
        e.set('typeface',KO)


def title(s, ref, text):
    for ph in ref.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE,PP_PLACEHOLDER.CENTER_TITLE):
            s.shapes._spTree.append(copy.deepcopy(ph._element)); break
    for ph in s.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE,PP_PLACEHOLDER.CENTER_TITLE):
            ph.text_frame.text=text
            for para in ph.text_frame.paragraphs:
                for r in para.runs: r.font.bold=True; r.font.color.rgb=BLUE
            break


def card(s, head, items, x, y, w, h, hcolor):
    box=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb=RGBColor(0xF3,0xF6,0xFB)
    box.line.color.rgb=hcolor; box.line.width=Pt(1.5); box.shadow.inherit=False
    tb=s.shapes.add_textbox(Inches(x+0.15),Inches(y+0.12),Inches(w-0.3),Inches(0.4))
    _kf(tb.text_frame.paragraphs[0].add_run(),14,True,hcolor); tb.text_frame.paragraphs[0].runs[0].text=head
    body=s.shapes.add_textbox(Inches(x+0.15),Inches(y+0.6),Inches(w-0.3),Inches(h-0.7))
    body.text_frame.word_wrap=True; first=True
    for it in items:
        txt,lvl=(it if isinstance(it,tuple) else (it,0)); lvl=min(lvl,2)
        p=body.text_frame.paragraphs[0] if first else body.text_frame.add_paragraph(); first=False
        p.space_after=Pt(3); p.level=lvl
        r=p.add_run(); r.text=("   "*lvl)+_B[lvl]+txt; _kf(r,_S[lvl],lvl==0,_C[lvl])


def note(s,text,top,color=ORANGE,size=12):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(top),Inches(12.3),Inches(0.6)); tb.text_frame.word_wrap=True
    _kf(tb.text_frame.paragraphs[0].add_run(),size,True,color); tb.text_frame.paragraphs[0].runs[0].text=text


def main():
    ap=argparse.ArgumentParser(); ap.add_argument("-f","--file",default="/home/autolab/Downloads/VLA_교육자료.pptx")
    a=ap.parse_args(); prs=Presentation(a.file)
    anchor=None
    for i,s in enumerate(prs.slides):
        if s.shapes and s.shapes[0].has_text_frame and "사전 학습 산출물" in s.shapes[0].text_frame.text:
            anchor=i; break
    if anchor is None:
        for i,s in enumerate(prs.slides):
            if s.shapes and s.shapes[0].has_text_frame and "사전 준비물" in s.shapes[0].text_frame.text:
                anchor=i; break
    ref=None
    for s in prs.slides:
        if s.shapes and s.shapes[0].has_text_frame and "핵심 원리" in s.shapes[0].text_frame.text:
            ref=s; break
    ref=ref or prs.slides[anchor]

    s=prs.slides.add_slide(ref.slide_layout)
    title(s, ref, "노트북이 없다면 — 무료 브라우저로 참여")
    # 상단 안내
    tb=s.shapes.add_textbox(Inches(0.5),Inches(1.25),Inches(12.3),Inches(0.5)); tb.text_frame.word_wrap=True
    _kf(tb.text_frame.paragraphs[0].add_run(),14,True,BLACK)
    tb.text_frame.paragraphs[0].runs[0].text="메인 트랙 = NVIDIA GPU 노트북(로컬 실습).  노트북이 없으면 아래 무료 브라우저 환경을 이용."

    card(s,"① Google Colab (무료 T4 GPU)", [
        "VLA(Qwen3-VL-2B) 추론·학습에 필요한 GPU 무료 제공",
        "산출물은 HuggingFace Hub에서 import (다운로드 자동)",
        "Gazebo 화면은 Xvfb+noVNC로 브라우저에서 관찰",
        "제공 노트북: colab/VLA_Simulation_Colab.ipynb",
    ], 0.4, 1.9, 6.15, 3.5, BLUE)

    card(s,"② The Construct (브라우저 ROS·Gazebo)", [
        "설치 불필요 · 무료 학생 계정 · 어떤 컴퓨터든",
        "Gazebo가 브라우저에 내장 → 시각화 쉬움",
        "ROSject 링크로 실습 환경 즉시 공유",
        "GPU 없음 → 고전/CNN 주행·시각화에 적합",
    ], 6.78, 1.9, 6.15, 3.5, GREEN)

    note(s,"GPU가 필요한 VLA 추론 = Colab(①),  Gazebo 시각화가 쉬운 곳 = The Construct(②) — 둘을 조합.",top=5.55)
    note(s,"⚠️ 이 프로젝트는 Gazebo Classic(2025.1 EOL) 기반 → 원본 재현은 Colab에서 Classic 그대로 구동이 정확.",
         top=6.15, color=GREY)

    lst=prs.slides._sldIdLst; ids=list(lst); newid=ids[-1]; anchor_id=ids[anchor]
    lst.remove(newid); anchor_id.addnext(newid)
    prs.save(a.file)
    print("inserted browser-track slide |", len(prs.slides._sldIdLst), "slides")


if __name__=="__main__":
    main()
