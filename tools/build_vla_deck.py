#!/usr/bin/env python3
"""
build_vla_deck.py — 연구실 템플릿(.pptx) 양식으로 VLA 교육 발표자료를 생성.

사용법:
    python3 tools/build_vla_deck.py \
        --template "/home/autolab/Downloads/연구실템플릿 .pptx" \
        -o "/home/autolab/Downloads/VLA_교육자료.pptx"
"""
import argparse, os, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 색상 (템플릿 테마 기반) ────────────────────────────────────────────────
BLUE   = RGBColor(0x00, 0x70, 0xC0)   # 제목색
ACC    = RGBColor(0x44, 0x72, 0xC4)   # accent1
DARK   = RGBColor(0x22, 0x28, 0x30)   # 코드박스 배경
CODEBG = RGBColor(0x1E, 0x1E, 0x1E)
CODEFG = RGBColor(0xE6, 0xE6, 0xE6)
GREEN  = RGBColor(0x70, 0xAD, 0x47)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
GREY   = RGBColor(0x59, 0x59, 0x59)
BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
FONT   = "Arial"
MONO   = "Consolas"

SW, SH = Inches(13.33), Inches(7.5)


def _layout(prs, name):
    for m in prs.slide_masters:
        for l in m.slide_layouts:
            if l.name == name:
                return l
    return prs.slide_layouts[0]


def _clear_slides(prs):
    """기존 슬라이드를 sldIdLst 항목 + 실제 파트까지 완전 제거 (파일명 충돌 방지)."""
    part = prs.part
    sldIdLst = prs.slides._sldIdLst
    for sid in list(sldIdLst):
        rId = sid.get(qn('r:id'))
        slide_part = part.related_part(rId)
        # 슬라이드가 참조하는 이미지 등은 남아도 무방; 슬라이드 파트만 드롭
        part.drop_rel(rId)
        try:
            prs.part.package._parts  # noop, keep ref
        except Exception:
            pass
        sldIdLst.remove(sid)


def _set(run, size=None, bold=None, color=None, mono=False, italic=None, font=None):
    """폰트 서식 강제. font를 명시하지 않으면 템플릿/테마 폰트를 상속(설정 안 함)."""
    if mono:            run.font.name = MONO
    elif font:          run.font.name = font
    if size is not None:  run.font.size = Pt(size)
    if bold is not None:  run.font.bold = bold
    if italic is not None: run.font.italic = italic
    if color is not None: run.font.color.rgb = color


def _ph(s):
    return {p.placeholder_format.idx: p for p in s.placeholders}


def add_title_slide(prs, title, subtitle):
    """표지 — 템플릿 TITLE 레이아웃 서식(54pt bold 등) 그대로 상속."""
    s = prs.slides.add_slide(_layout(prs, "TITLE"))
    ph = _ph(s)
    if 0 in ph: ph[0].text_frame.text = title
    if 1 in ph: ph[1].text_frame.text = subtitle
    return s


def add_section(prs, title, subtitle=None):
    """섹션 구분 — 템플릿 SECTION_HEADER 서식(60pt) 상속."""
    s = prs.slides.add_slide(_layout(prs, "SECTION_HEADER"))
    ph = _ph(s)
    if 0 in ph: ph[0].text_frame.text = title
    if subtitle and 1 in ph: ph[1].text_frame.text = subtitle
    return s


def add_content(prs, title):
    """본문 — 템플릿 OBJECT 레이아웃(제목 32pt + ❑/🔾/◉ 개요) 상속."""
    s = prs.slides.add_slide(_layout(prs, "OBJECT"))
    ph = _ph(s)
    if 0 in ph: ph[0].text_frame.text = title
    if 1 in ph: ph[1].text_frame.clear()   # 본문 placeholder 준비 (개요 채움)
    return s


def body(s, items):
    """OBJECT 레이아웃의 BODY placeholder에 개요를 채운다.
    → 템플릿 고유의 불릿(❑/🔾/◉)·들여쓰기·폰트·크기를 그대로 상속.
    items: [(text, level), ...]  level 0/1/2. 빈 문자열은 빈 줄."""
    ph = _ph(s).get(1)
    if ph is None:
        return bullets(s, items)  # 폴백
    tf = ph.text_frame; tf.word_wrap = True
    items = [it for it in items if (it[0] if isinstance(it, tuple) else it).strip()]  # 빈 줄 제거(불릿 잔상 방지)
    first = True
    for it in items:
        txt, lvl = (it if isinstance(it, tuple) else (it, 0))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.add_run().text = txt     # 폰트/불릿/색은 템플릿 lstStyle 상속 (강제 안 함)
    return ph


def body_box(s, items, left, top, width, height):
    """BODY 개요를 특정 위치/크기로 배치해야 할 때 (코드박스와 공존하는 슬라이드)."""
    ph = _ph(s).get(1)
    if ph is not None:
        ph.left, ph.top, ph.width, ph.height = (Inches(left), Inches(top),
                                                 Inches(width), Inches(height))
        return body(s, items)
    return bullets(s, items, left, top, width, height)


def bullets(s, items, left=0.55, top=1.35, width=12.2, height=5.3, size=16):
    """폴백용 수동 텍스트박스 (placeholder가 없을 때만 사용)."""
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        txt, lvl = (it if isinstance(it, tuple) else (it, 0))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl; p.space_after = Pt(5)
        run = p.add_run(); run.text = txt
        _set(run, size - lvl, None, BLACK if lvl == 0 else GREY)
    return tb


def code_box(s, lines, left=0.55, top=None, width=12.2, height=None, size=12.5, title=None):
    """어두운 배경의 코드/명령어 박스."""
    if top is None: top = 4.4
    if height is None: height = 0.5 + 0.28*len(lines)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb = CODEBG
    box.line.color.rgb = RGBColor(0x40,0x40,0x40); box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(1); p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = ln if ln else " "
        # 주석(#) 회색, $ 프롬프트 초록, 나머지 밝은 회색
        if ln.strip().startswith("#"):
            _set(r, size, False, RGBColor(0x8A,0xA8,0x8A), mono=True)
        else:
            _set(r, size, False, CODEFG, mono=True)
            # 앞의 $ 만 초록으로: 간단히 전체 밝게 처리
    return box


def note(s, text, left=0.55, top=6.55, width=12.2, color=ORANGE, size=12):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    _set(r, size, True, color)
    return tb


def hdr(s, text, top, left=0.5, size=16, color=None):
    """코드박스와 함께 쓰는 한 줄 소제목 (테마 폰트 상속, 굵게)."""
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(12.3), Inches(0.45))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    _set(r, size, True, color if color else BLUE)
    return tb


def label(s, text, left, top, width, height, size=13, color=BLACK, bold=True,
          fill=None, align=PP_ALIGN.CENTER):
    if fill is not None:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.fill.solid(); box.fill.fore_color.rgb = fill
        box.line.fill.background(); box.shadow.inherit = False
        tf = box.text_frame
    else:
        box = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        _set(r, size, bold, color)
    return box


def picture(s, path, left, top, width=None, height=None):
    if not os.path.exists(path): return None
    kw = {}
    if width: kw["width"] = Inches(width)
    if height: kw["height"] = Inches(height)
    return s.shapes.add_picture(path, Inches(left), Inches(top), **kw)


def arrow(s, l, t, w, h=0.0, color=ACC):
    """단순 오른쪽/아래 화살표."""
    shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW if h == 0 else MSO_SHAPE.DOWN_ARROW,
                             Inches(l), Inches(t), Inches(w if w else 0.5), Inches(h if h else 0.4))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def conn(s, l, t, w, color=GREY, size=10):
    """토픽 라벨(화살표 위 작은 글씨)."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "▶"
    _set(r, size, False, color)
    return tb


def table(s, rows, left, top, width, col_widths=None, header_fill=ACC,
          head_size=12, body_size=11, row_h=0.42):
    """rows[0]=헤더. 셀은 (텍스트) 또는 (텍스트,색). 폰트는 지정."""
    nrows, ncols = len(rows), len(rows[0])
    gt = s.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                            Inches(width), Inches(row_h*nrows)).table
    gt.first_row = False; gt.horz_banding = True
    if col_widths:
        for i, w in enumerate(col_widths):
            gt.columns[i].width = Inches(w)
    for r_i, row in enumerate(rows):
        gt.rows[r_i].height = Inches(row_h if r_i else row_h*0.9)
        for c_i, cell in enumerate(row):
            txt, col = (cell if isinstance(cell, tuple) else (cell, None))
            c = gt.cell(r_i, c_i)
            c.margin_left = Inches(0.06); c.margin_right = Inches(0.04)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
            else:
                c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF) if r_i%2 else RGBColor(0xF0,0xF4,0xFA)
            tf = c.text_frame; tf.word_wrap = True
            first = True
            for line in txt.split("\n"):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                rr = p.add_run(); rr.text = line
                is_head = (r_i == 0)
                mono = ("_" in line or "/" in line) and not is_head and c_i <= 2
                _set(rr, head_size if is_head else body_size, is_head,
                     RGBColor(0xFF,0xFF,0xFF) if is_head else (col or BLACK),
                     mono=mono)
    return gt


# ═══════════════════════════════════════════════════════════════════════════
def build(prs):
    WS = "~/ros2_autonomous_vehicle_simulation"
    IMG = "/home/autolab"

    # 0. 표지
    add_title_slide(prs,
        "VLA 자율주행 A to Z",
        "Vision-Language-Action 기반 자율주행 시뮬레이션\n"
        "데이터 수집 · 학습 · 추론 · 실전 완전 가이드\n\n자동화연구실 · 성균관대학교")

    # 1. 목차
    s = add_content(prs, "목차")
    body(s, [
        ("1. VLA란 무엇인가 — 개념과 원리", 0),
        ("2. 시스템 아키텍처 (Fast / Slow 이중 시스템)", 0),
        ("3. 환경 준비 — 설치 · 빌드 · HuggingFace · GPU", 0),
        ("4. 데이터 수집 — 수동주행부터 라벨 생성까지", 0),
        ("5. 학습 — 무엇을, 어떻게 배우나", 0),
        ("6. 추론 & 자율주행 — 모델 로드부터 평가까지", 0),
        ("7. 부록 — 트랙 GT 시각화 도구", 0),
        ("8. 전체 흐름 요약 · 트러블슈팅", 0),
    ])

    # ── SECTION 1 ──────────────────────────────────────────────────────────
    add_section(prs, "1. VLA란 무엇인가")

    s = add_content(prs, "VLA = Vision · Language · Action")
    body(s, [
        ("VLA는 세 가지를 하나의 모델로 잇는 자율주행 패러다임이다:", 0),
        ("Vision (시각) — 카메라 이미지로 주변을 인식", 1),
        ("Language (언어) — '1차선으로 한 바퀴 돌아' 같은 자연어 명령 이해", 1),
        ("Action (행동) — 곧바로 조향·속도 같은 주행 행동 출력", 1),
        ("기존 방식과의 차이", 0),
        ("고전 파이프라인 = 차선검출 → 경로계획 → 제어 (모듈마다 규칙을 사람이 설계)", 1),
        ("VLA = 이미지+명령 → (신경망 한 번) → 행동. 사람의 주행을 모방 학습", 1),
        ("장점: 처음 보는 상황·자연어 지시에 유연 / 단점: 데이터·연산 필요", 1),
    ])

    s = add_content(prs, "이 프로젝트의 VLA — 한 줄 요약")
    label(s, "카메라\n이미지", 0.5, 2.7, 1.7, 1.1, 15, RGBColor(0xFF,0xFF,0xFF), True, ACC)
    arrow(s, 2.3, 3.0, 0.5)
    label(s, "Qwen3-VL\n비전 인코더\n(고정)", 2.9, 2.6, 2.0, 1.3, 14, RGBColor(0xFF,0xFF,0xFF), True, GREY)
    arrow(s, 5.0, 3.0, 0.5)
    label(s, "작은 Head\n(학습)\n→ 6 웨이포인트", 5.6, 2.55, 2.2, 1.4, 13, RGBColor(0xFF,0xFF,0xFF), True, GREEN)
    arrow(s, 7.9, 3.0, 0.5)
    label(s, "Pure-Pursuit\n(고전 제어)", 8.5, 2.7, 2.0, 1.1, 13, BLACK, True, RGBColor(0xFF,0xE0,0x99))
    arrow(s, 10.6, 3.0, 0.5)
    label(s, "조향/속도\nMotionCommand", 11.2, 2.6, 1.9, 1.3, 12, RGBColor(0xFF,0xFF,0xFF), True, ORANGE)
    body_box(s, [
        ("카메라 이미지가 고정된 Qwen3-VL 비전 인코더를 거쳐, 작은 학습 헤드가", 0),
        ("'앞으로 지나갈 6개 웨이포인트'를 예측 → 고전 제어기가 조향으로 변환한다.", 0),
    ], 0.42, 4.3, 12.4, 1.6)
    note(s, "핵심: 20억 파라미터 모델 전체가 아니라, 그 위에 얹은 '작은 부분'만 학습한다.")

    s = add_content(prs, "핵심 원리 3가지 (왜 이렇게 설계했나)")
    body(s, [
        ("① 비전은 고정(freeze), 작은 부분만 학습 — 그래서 LoRA", 0),
        ("Qwen3-VL-2B(20억)를 처음부터 학습은 불가/불필요. 비전 인코더는 고정하고", 1),
        ("작은 Head(수 MB)만 학습. 더 필요할 때만 LoRA(저차원 어댑터)로 극소수만 학습", 1),
        ("② 행동을 '웨이포인트'로 예측 (조향 글자 X)", 0),
        ("'D -3 40' 같은 글자 출력은 느리고 부정확 → 앞으로 갈 6개 점(ex,ey) 회귀 예측", 1),
        ("고전 pure-pursuit이 그 점들을 따라가며 조향 계산 (빠르고 안정적)", 1),
        ("③ 공간 정보를 뭉개지 않는다 (spatial head)", 0),
        ("70개 토큰을 평균내면 '차선 내 좌우 위치'가 사라져 직진이 흔들림", 1),
        ("토큰별 위치를 보존하는 Head로 교체 → 오차(wpMAE) 0.56 → 0.284로 개선", 1),
    ])

    # ── SECTION 2 ──────────────────────────────────────────────────────────
    add_section(prs, "2. 시스템 아키텍처")

    s = add_content(prs, "Fast / Slow 이중 시스템")
    label(s, "FAST — 실시간 주행", 0.6, 1.5, 5.9, 0.5, 16, RGBColor(0xFF,0xFF,0xFF), True, GREEN)
    label(s, "vla_lora_drive_node\n"
             "· 매 프레임 카메라 → 웨이포인트 → 조향\n"
             "· 비전 1-pass ~10ms (98FPS), 실시간\n"
             "· 차선 추종을 담당",
          0.6, 2.1, 5.9, 1.9, 14, BLACK, False, RGBColor(0xEE,0xF6,0xEC), align=PP_ALIGN.LEFT)
    label(s, "SLOW — 지능/명령 해석", 6.8, 1.5, 5.9, 0.5, 16, RGBColor(0xFF,0xFF,0xFF), True, ORANGE)
    label(s, "vla_brain_node (Qwen3-VL-2B)\n"
             "· 자유로운 자연어 명령 해석\n"
             "· 처음 보는 물체·신호등 zero-shot 판단\n"
             "· 표준 명령으로 바꿔 vla/command 발행",
          6.8, 2.1, 5.9, 1.9, 14, BLACK, False, RGBColor(0xFD,0xF0,0xE6), align=PP_ALIGN.LEFT)
    label(s, "안전 로직 (결정론적 코드) — 라이다 전방 정지·기하 계산 등 '신뢰가 중요한 부분'은\n"
             "LLM이 아니라 규칙 기반 코드가 담당한다.", 0.7, 4.3, 12.0, 1.0, 15, BLACK, True,
          RGBColor(0xED,0xED,0xED))
    note(s, "역할 분담: 빠른 반응(Fast) + 똑똑한 해석(Slow) + 확실한 안전(코드)", top=5.7)

    # ── 노드 흐름도 (닫힌 루프 강조) ──
    s = add_content(prs, "노드 데이터 흐름도 — 닫힌 루프")
    def nbox(text, l, t, w, h, fill, fg=RGBColor(0xFF,0xFF,0xFF), sz=12):
        label(s, text, l, t, w, h, sz, fg, True, fill)
    def tlabel(text, l, t, w=1.4, color=GREY, sz=9):
        label(s, text, l, t, w, 0.28, sz, color, False, align=PP_ALIGN.CENTER)

    # ① 명령 계통 (상단, 좌→우)
    label(s, "① 명령 (Command)", 0.5, 1.2, 3.0, 0.35, 12, RGBColor(0x7A,0x5A,0x9A), True, align=PP_ALIGN.LEFT)
    nbox("vla_gui\n명령 콘솔(GUI)", 0.5, 1.65, 2.4, 0.95, RGBColor(0x7A,0x5A,0x9A))
    nbox("vla_brain_node\nQwen3-VL 브레인", 3.55, 1.65, 2.7, 0.95, ORANGE)
    nbox("vla_lora_drive_node\nVLA 실시간 드라이버", 7.0, 1.65, 3.1, 0.95, GREEN)
    arrow(s, 2.95, 2.0, 0.55); tlabel("/nl_command", 2.55, 1.4, 1.4, RGBColor(0x7A,0x5A,0x9A))
    arrow(s, 6.3, 2.0, 0.6); tlabel("vla/command", 6.05, 1.4, 1.4, ORANGE)

    # ② 제어 계통 (우측, 아래로)
    label(s, "② 제어 (Control)", 7.0, 2.75, 3.0, 0.32, 12, GREY, True, align=PP_ALIGN.LEFT)
    arrow(s, 8.15, 3.1, 0.55, h=0.75)
    tlabel("topic_control_signal\n(MotionCommand)", 8.75, 3.15, 2.3, GREY, 9)
    nbox("simulation_sender_node\nMotionCommand → cmd_vel", 7.0, 3.95, 3.1, 0.9, GREY)
    nbox("Gazebo 시뮬레이터\n차량·카메라·라이다·odom", 10.55, 3.9, 2.4, 1.0, ACC)
    arrow(s, 10.15, 4.3, 0.4); tlabel("/cmd_vel", 9.5, 3.95, 1.0, GREY)

    # ③ 센서 피드백 (하단 리턴 레일 — 폐루프 강조)
    rail = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(0.7), Inches(5.75),
                              Inches(11.1), Inches(0.55))
    rail.fill.solid(); rail.fill.fore_color.rgb = ACC
    rail.line.fill.background(); rail.shadow.inherit = False
    tf = rail.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "③ 센서 피드백 (폐루프):  camera/image_raw  ·  scan  ·  /odom  →  다시 드라이버/브레인으로"
    _set(r, 12, True, RGBColor(0xFF,0xFF,0xFF))
    # Gazebo ↓ 레일,  레일 ↑ 드라이버(및 브레인)
    down = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(11.55), Inches(4.9), Inches(0.35), Inches(0.9))
    down.fill.solid(); down.fill.fore_color.rgb = ACC; down.line.fill.background(); down.shadow.inherit = False
    up = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(4.6), Inches(2.75), Inches(0.35), Inches(3.05))
    up.fill.solid(); up.fill.fore_color.rgb = ACC; up.line.fill.background(); up.shadow.inherit = False
    label(s, "이미지·라이다·위치\n피드백", 5.05, 3.6, 2.0, 0.7, 10, ACC, True, align=PP_ALIGN.LEFT)

    # 브레인↔드라이버 동기 (up 화살표 왼쪽에 배치)
    tlabel("vla/cur_lane\n(차선 동기)", 3.35, 2.75, 1.15, GREY, 9)
    note(s, "명령 → 제어 → 센서 피드백이 초당 수십 번 순환하는 닫힌 루프(closed loop)로 차가 스스로 달린다.",
         top=6.55)

    # ── 노드별 상세 표 ──
    s = add_content(prs, "노드별 상세 — 구독 · 발행 · 역할")
    table(s, [
        ["노드", "구독 (입력)", "발행 (출력)", "역할"],
        [("vla_lora_drive_node", GREEN),
         "camera/image_raw\nscan, /odom\nvla/command",
         "topic_control_signal\nvla/cur_lane",
         "★ 실시간 드라이버. 비전→웨이포인트→pure-pursuit 조향, 라이다 정지/회피, 랩 카운트"],
        [("vla_brain_node", ORANGE),
         "camera/image_raw\nnl_command\nvla/cur_lane",
         "vla/command",
         "Qwen3-VL 브레인. 자연어→표준명령 변환, 신호등 zero-shot 감시(일시정지/재개)"],
        [("vla_gui", RGBColor(0x7A,0x5A,0x9A)),
         "vla/command",
         "nl_command\nvla/command",
         "명령 콘솔(PySide6). 자연어 입력→브레인, 퀵버튼→드라이버 직접"],
        [("simulation_sender_node", GREY),
         "topic_control_signal",
         "/cmd_vel",
         "MotionCommand(조향±7·속도)→Ackermann Twist 변환 브리지"],
        [("load_ego_car_node", ACC),
         "— (1회 스폰)",
         "—",
         "prius_hybrid 모델을 ego_vehicle 로 Gazebo에 스폰"],
        [("Gazebo 플러그인", ACC),
         "/cmd_vel",
         "camera/image_raw\nscan, /odom",
         "차량 물리·센서. 카메라/라이다/odom 발행, cmd_vel로 주행"],
    ], 0.35, 1.2, 12.6,
       col_widths=[2.5, 2.5, 2.4, 5.2], head_size=11.5, body_size=10, row_h=0.72)
    note(s, "핵심 제어 메시지 interfaces_pkg/msg/MotionCommand: int32 steering(-7~7), left_speed, right_speed",
         top=6.5, size=11)

    # ── 데이터 순환 (제어 체인) ──
    s = add_content(prs, "제어 신호 한 바퀴 — MotionCommand의 여정")
    code_box(s, [
        "① vla_lora_drive_node : 카메라·라이다·odom·명령을 보고 매 프레임 판단",
        "        └─ MotionCommand(steering, left_speed, right_speed)",
        "                       │  topic_control_signal",
        "                       ▼",
        "② simulation_sender_node : 조향±7·속도를 Ackermann 자전거 모델로 변환",
        "        └─ geometry_msgs/Twist (angular.z=조향, linear.x=속도)",
        "                       │  /cmd_vel",
        "                       ▼",
        "③ Gazebo Ackermann 플러그인 : 바퀴를 굴려 차량 이동",
        "        └─ 카메라/라이다/odom 재발행  →  다시 ①로 (폐루프)",
    ], top=1.5, size=13, height=4.2)
    note(s, "이 폐루프가 초당 수십 번 반복되며 차가 스스로 트랙을 돈다.")

    # ── SECTION 3 ──────────────────────────────────────────────────────────
    add_section(prs, "3. 환경 준비")

    s = add_content(prs, "워크스페이스 구조")
    code_box(s, [
        "~/ros2_autonomous_vehicle_simulation/",
        " ├─ src/",
        " │   ├─ interfaces_pkg/        # 커스텀 메시지 (MotionCommand 등)",
        " │   ├─ camera_perception_pkg/ # 카메라·YOLOv8·차선·신호등",
        " │   ├─ lidar_perception_pkg/  # 라이다",
        " │   ├─ decision_making_pkg/   # 경로·모션 계획 (고전)",
        " │   ├─ simulation_pkg/        # Gazebo world·모델·launch",
        " │   ├─ qwen_vl_pkg/           # Qwen3-VL 상위 VLA 노드",
        " │   └─ gui_pkg/               # GUI·GT 어노테이터",
        " ├─ lora_pipeline/             # ★ VLA 학습·주행 스크립트 (핵심)",
        " └─ tools/                     # GT 시각화 도구",
    ], top=1.5, size=13.5, height=4.6)
    note(s, "VLA 실습의 핵심 스크립트는 모두 lora_pipeline/ 안에 있다 (ROS 패키지가 아닌 독립 스크립트).")

    s = add_content(prs, "설치 & 빌드")
    hdr(s, "① 의존성 설치 (최초 1회)", 1.4)
    code_box(s, [
        "cd ~/ros2_autonomous_vehicle_simulation",
        "sh install.sh                         # gazebo, ultralytics, transformers 등",
        "pip install torch transformers peft accelerate   # VLA 추가 의존성",
    ], top=1.85, size=13, height=1.2)
    hdr(s, "② 빌드 (interfaces_pkg를 반드시 먼저)", 3.3)
    code_box(s, [
        "source /opt/ros/humble/setup.bash",
        "colcon build --packages-select interfaces_pkg --allow-overriding interfaces_pkg",
        "source install/local_setup.bash",
        "colcon build --symlink-install --packages-select \\",
        "  camera_perception_pkg decision_making_pkg debug_pkg simulation_pkg \\",
        "  lidar_perception_pkg mission_control_pkg gui_pkg qwen_vl_pkg",
        "source install/local_setup.bash",
    ], top=3.8, size=12.5, height=2.4)

    s = add_content(prs, "HuggingFace — 모델 다운로드")
    body_box(s, [
        ("이 프로젝트가 쓰는 모델: Qwen/Qwen3-VL-2B-Instruct", 0),
        ("코드가 from_pretrained 로 자동 다운로드 (최초 실행 시 ~/.cache/huggingface 로)", 1),
        ("미리 받아두면(권장) 첫 실행이 끊기지 않는다:", 0),
    ], 0.42, 1.35, 12.4, 1.6)
    code_box(s, [
        "pip install -U huggingface_hub",
        "# 방법1) CLI 로 미리 다운로드",
        "huggingface-cli download Qwen/Qwen3-VL-2B-Instruct",
        "",
        "# 방법2) 캐시 위치 지정 (용량 큰 디스크로)",
        "export HF_HOME=/data/hf_cache        # ~/.bashrc 에 추가 권장",
        "",
        "# 방법3) 게이트 모델/속도 — 로그인 & 고속 다운로드",
        "huggingface-cli login                # 토큰 입력 (필요시)",
        "pip install hf_transfer; export HF_HUB_ENABLE_HF_TRANSFER=1",
    ], top=3.0, size=12.5, height=3.0)
    note(s, "오프라인 실행: 한번 받은 뒤 export HF_HUB_OFFLINE=1 로 네트워크 없이 로드 가능", top=6.5)

    s = add_content(prs, "GPU 가속 · 더 빠르게 쓰는 법")
    hdr(s, "신경망은 GPU에서 실행 (device_map='cuda:0', bfloat16, SDPA attention)", 1.35)
    code_box(s, [
        "nvidia-smi                                   # GPU 상태 확인",
        "python3 -c \"import torch; print(torch.cuda.is_available(), torch.cuda.get_device_name(0))\"",
    ], top=1.85, size=12.5, height=0.95)
    body_box(s, [
        ("더 빠르게 만드는 실전 팁", 0),
        ("Fast 경로는 '토큰 생성' 없이 비전 1-pass만 → ~10ms(98FPS). 글자 생성 대비 수십 배 빠름", 1),
        ("학습 시 비전 특징을 캐시(dataset/*_cache.pt) → train_head_fast는 에폭당 수 초", 1),
        ("입력 해상도 축소(PIXELS=320x240)와 bfloat16으로 연산량 ↓", 1),
        ("추론만 필요하면 vla_lora_head_fast.pt(비전 고정) 사용 — LoRA보다 가벼움", 1),
        ("Gazebo 카메라는 의도적으로 소프트웨어 렌더 유지 (PRIME offload는 카메라 readback이 느림)", 1),
        ("즉, '무거운 계산=GPU(torch), Gazebo 렌더=CPU' 로 역할을 나눈 것이 정답 설정", 2),
    ], 0.42, 3.05, 12.4, 3.4)

    # ── SECTION 4 ──────────────────────────────────────────────────────────
    add_section(prs, "4. 데이터 수집")

    s = add_content(prs, "데이터 수집 개요")
    body_box(s, [
        ("목표: '사람이 어떻게 운전했는가'를 모아 모델의 정답표(라벨)를 만든다.", 0),
        ("4단계 흐름", 0),
        ("[1] teleop_sim 실행 — 차+카메라+제어브리지만 (플래너 없음)", 1),
        ("[2] teleop_keyboard.py — 키보드로 직접 운전", 1),
        ("[3] manual_collect.py — 이미지+조작+위치 저장", 1),
        ("[4] build_wp_from_manual.py — 웨이포인트 라벨 생성", 1),
        ("각 단계는 별도 터미널. 모든 새 터미널에서 먼저 환경 소스:", 0),
    ], 0.42, 1.35, 12.4, 4.0)
    code_box(s, [
        "cd ~/ros2_autonomous_vehicle_simulation",
        "source /opt/ros/humble/setup.bash && source install/local_setup.bash",
        "export DISPLAY=:1        # 헤드리스 서버인 경우",
    ], top=5.5, size=12.5, height=1.2)

    s = add_content(prs, "[1]-[2] 수동 주행 & 키보드 조작")
    hdr(s, "터미널 A — 최소 시뮬레이터 실행", 1.35)
    code_box(s, [
        "sudo killall -9 gazebo gzserver gzclient 2>/dev/null",
        "ros2 launch simulation_pkg teleop_sim.launch.py",
    ], top=1.8, size=13, height=0.9)
    hdr(s, "터미널 B — 키보드 텔레옵 (이 창에 포커스가 있어야 키 입력됨)", 2.95)
    code_box(s, [
        "python3 lora_pipeline/teleop_keyboard.py     # 조향 반대면 --invert",
    ], top=3.4, size=13, height=0.6)
    label(s, "w / s : 속도 +20 / -20 (최대 ±255)\n"
             "a / d : 조향 왼쪽 / 오른쪽 (최대 ±7)\n"
             "space : 정지        x : 조향·속도 리셋        q : 종료",
          0.6, 4.2, 12.1, 1.4, 15, BLACK, True, RGBColor(0xED,0xF2,0xFB))
    note(s, "주의: teleop_sim 은 키보드 노드를 켜지 않는다 → teleop_keyboard.py 를 별도로 실행. "
             "'c'키(중앙정렬)는 미구현 → x 사용", top=5.9)

    s = add_content(prs, "[3] 주행 데이터 수집")
    hdr(s, "터미널 C — 사람이 운전하는 동안 저장 (--lane: 0=1차선, 1=2차선)", 1.35)
    code_box(s, [
        "python3 lora_pipeline/manual_collect.py --lane 0",
    ], top=1.85, size=13, height=0.6)
    body_box(s, [
        ("저장 결과", 0),
        ("이미지 → lora_pipeline/manual_demos/images/man_L0_<시간>_<번호>.jpg", 1),
        ("라벨   → lora_pipeline/manual_demos/labels.csv", 1),
        ("컬럼: fname, steering, speed, x, y, yaw, lane  (조작 + 차 위치/방향)", 2),
        ("초당 10프레임, 차가 움직일 때만 기록 (정지 프레임 제외)", 1),
        ("좋은 데이터 팁", 0),
        ("차선 중앙을 부드럽게, 곡선 구간 충분히, 양쪽 차선(--lane 0/1) 모두 수집", 1),
        ("많고 다양할수록 모델이 잘 배운다. Ctrl+C 로 종료", 1),
    ], 0.42, 2.7, 12.4, 3.7)

    s = add_content(prs, "[4] 웨이포인트 라벨 생성")
    body_box(s, [
        ("수집한 (이미지+위치)를 '앞으로 갈 6개 점' 라벨로 변환", 0),
        ("트랙 중심선 GT 파일이 필요 (~/track_gt_manual.json, ~/track_gt_outward_centerline.json)", 1),
    ], 0.42, 1.35, 12.4, 1.2)
    code_box(s, [
        "python3 lora_pipeline/build_wp_from_manual.py",
        "# 출력: lora_pipeline/manual_wp_labels.csv",
        "# 컬럼: path, ex0,ey0, ex1,ey1, ... ex5,ey5, lane   (차량좌표계: ex=전방, ey=좌측, m)",
    ], top=2.7, size=12.5, height=1.3)
    bullets(s, [
        ("의미: 각 장면에서 차의 현재 위치·방향 기준으로 중심선을 따라", 0),
        ("앞쪽 6개 지점을 투영 = '이 장면에서 사람은 이런 궤적으로 갔다'는 정답표", 0),
    ], left=0.55, top=4.3, width=12.2, size=15)
    note(s, "이 GT 중심선을 눈으로 확인하는 도구는 뒤의 '부록: GT 시각화'에서 다룬다.", top=5.6)

    # ── SECTION 5 ──────────────────────────────────────────────────────────
    add_section(prs, "5. 학습")

    s = add_content(prs, "무엇을, 어떻게 배우나")
    body(s, [
        ("최종 목표: 이미지 → 6개 웨이포인트 매핑을 학습", 0),
        ("① 이미지 → 특징: 고정된 Qwen3-VL 비전이 70토큰 × 2048차원 생성", 0),
        ("이 토큰들을 평균내지 않고(spatial) 그대로 사용 → 좌우 위치 보존", 1),
        ("한 번 계산한 특징은 캐시(dataset/*_cache.pt) → 재학습이 매우 빠름", 1),
        ("② Head가 토큰 → 6웨이포인트를 회귀 학습", 0),
        ("두 가지 방식", 0),
        ("(A) train_head_fast — 비전 고정 + Head만 (빠름, 기본/권장)", 1),
        ("(B) train_vla_lora — 비전 LoRA + Head 동시 (정밀, 오래 걸림)", 1),
    ])

    s = add_content(prs, "Head 구조 (spatial head)")
    code_box(s, [
        "입력: 비전 토큰 (70, 2048)",
        "  → 토큰별 Linear(2048 → 64)          # 각 위치의 특징 압축",
        "  → + 위치 임베딩 (1, 70, 64)          # '어느 위치 토큰인지' 정보 부여",
        "  → flatten → Linear(70*64 → 512)      # 공간 정보를 한데 모음",
        "  → FiLM 차선 조건 Embedding(4, 1024)  # 1차선/2차선·차선유지/변경 조건",
        "  → MLP(512 → 256 → 12)                # 6점 x (ex, ey) = 12 출력",
        "출력: 6개 웨이포인트 (ex, ey)",
    ], top=1.6, size=13, height=2.9)
    body_box(s, [
        ("핵심 포인트: '평균 풀링' 대신 '토큰별 + 위치임베딩'을 쓰는 것이 성능의 열쇠", 0),
        ("FiLM = 차선 번호에 따라 특징을 조절(γ,β) → 같은 이미지라도 목표 차선에 맞게 예측", 0),
    ], 0.42, 4.75, 12.4, 1.6)

    s = add_content(prs, "[5] 모델 학습 — 명령어")
    hdr(s, "(A) 빠른 학습 — 비전 고정 + Head만  [권장, 처음엔 이것부터]", 1.35)
    code_box(s, [
        "python3 lora_pipeline/train_head_fast.py            # 기본",
        "EPOCHS=60 python3 lora_pipeline/train_head_fast.py  # 에폭 조절",
        "# 출력: lora_pipeline/vla_lora_head_fast.pt  ← 주행 노드가 기본으로 로드",
    ], top=1.8, size=12.5, height=1.35)
    hdr(s, "(B) 정밀 학습 — 비전 LoRA + Head 동시  [더 오래 걸림]", 3.45)
    code_box(s, [
        "python3 lora_pipeline/train_vla_lora.py",
        "# 출력: lora_pipeline/vla_lora_adapter/ (LoRA 어댑터) + vla_lora_head.pt",
        "# 수동 데모(manual_wp_labels.csv)도 함께 사용 (USE_MANUAL=1 기본)",
    ], top=3.9, size=12.5, height=1.35)
    note(s, "GPU 권장. (A)는 특징 캐시 덕에 에폭당 수 초, (B)는 비전까지 미세조정해 수 분~십수 분", top=5.6)

    s = add_content(prs, "학습 산출물 정리")
    code_box(s, [
        "lora_pipeline/",
        " ├─ vla_lora_head_fast.pt     # (A) 비전고정+Head — 주행 노드 기본 로드 ★",
        " ├─ vla_lora_head.pt          # (B) LoRA 학습 시 함께 나오는 Head",
        " ├─ vla_lora_adapter/         # (B) 비전 LoRA 어댑터 (adapter_model.safetensors)",
        " └─ dataset/*_cache.pt        # 비전 특징 캐시 (재학습 가속용)",
    ], top=1.6, size=13, height=2.4)
    body_box(s, [
        ("기본 주행은 (A) vla_lora_head_fast.pt = '비전 고정 + spatial head'", 0),
        ("(B) LoRA 어댑터로 주행하려면 환경변수 VLA_USE_ADAPTER=1 로 명시적으로 켠다", 0),
    ], 0.42, 4.3, 12.4, 1.6)

    # ── SECTION 6 ──────────────────────────────────────────────────────────
    add_section(prs, "6. 추론 & 자율주행")

    s = add_content(prs, "추론 파이프라인 (매 프레임)")
    code_box(s, [
        "camera BGR 이미지",
        " → 프로세서 (320x240)",
        " → Qwen3-VL 비전 (1-pass, ~10ms)  →  토큰 (70, 2048)   [평균X]",
        " → Head  →  6 웨이포인트 (ex, ey) x 스케일",
        " → Pure-Pursuit:  전방 0.55m 이상 첫 점 목표",
        "      he = atan2(ey, ex);  steer = clamp(-he * 13,  -7 .. +7)",
        " → MotionCommand(steering, left_speed, right_speed)  →  topic_control_signal",
    ], top=1.6, size=13, height=3.0)
    body_box(s, [
        ("라이다(scan)로 전방 장애물 시 감속/정지 — 안전은 결정론적 코드가 담당", 0),
        ("추론은 '토큰 생성' 없이 비전 특징만 뽑는 1-pass라 실시간(수십~수백 FPS 급)", 0),
    ], 0.42, 4.75, 12.4, 1.6)

    s = add_content(prs, "[6] 학습한 모델로 자율주행")
    code_box(s, [
        "cd ~/ros2_autonomous_vehicle_simulation",
        "sudo killall -9 gazebo gzserver gzclient 2>/dev/null",
        "source install/setup.bash",
        "ros2 launch lora_pipeline/vla_drive.launch.py",
        "#   옵션:  brain:=false (Qwen 브레인 끄기)  gui:=false  gzclient:=true (3D창)",
    ], top=1.5, size=13, height=2.0)
    hdr(s, "LoRA 어댑터로 주행하려면 (선택)", 3.8)
    code_box(s, [
        "VLA_USE_ADAPTER=1 \\",
        "VLA_HEAD=$PWD/lora_pipeline/vla_lora_head.pt \\",
        "VLA_ADAPTER=$PWD/lora_pipeline/vla_lora_adapter \\",
        "ros2 launch lora_pipeline/vla_drive.launch.py",
    ], top=4.25, size=12.5, height=1.7)
    note(s, "이 launch 하나가 Gazebo·차량스폰·제어브리지·VLA 주행노드를 모두 띄운다.", top=6.2)

    s = add_content(prs, "[7]-[8] 자연어 명령 & 평가")
    hdr(s, "자연어 명령 보내기 (새 터미널)", 1.35)
    code_box(s, [
        "source install/local_setup.bash",
        "ros2 topic pub --once /vla/command std_msgs/String \"{data: '1차선 한바퀴 돌아'}\"",
        "#  '멈춰'  |  '2차선 한바퀴 돌아'  등",
    ], top=1.8, size=12.5, height=1.25)
    hdr(s, "성능 평가", 3.35)
    code_box(s, [
        "python3 lora_pipeline/eval_lap.py --lane 0     # 폐루프: 중심선 이탈/커버리지/지터",
        "python3 lora_pipeline/one_lap_test.py          # 한 바퀴 자동 완료 판정",
    ], top=3.8, size=12.5, height=1.0)
    note(s, "자유 문장은 vla_brain_node(Qwen)가 /nl_command 로 받아 표준 vla/command 로 변환한다.", top=5.2)

    # ── SECTION 7 ──────────────────────────────────────────────────────────
    add_section(prs, "7. 부록 — 트랙 GT 시각화")

    s = add_content(prs, "GT 데이터란 & 좌표 변환")
    body_box(s, [
        ("GT(Ground-Truth) = 트랙 차선 중심선의 '정답' 좌표. 웨이포인트 라벨 생성의 기준", 0),
        ("파일 위치(~/): track_gt_manual.json(픽셀+월드), track_gt_*_demo.json(월드)", 1),
        ("두 좌표계: 픽셀(이미지) ↔ 월드(미터). 변환식(gt_annotator.py):", 0),
    ], 0.42, 1.35, 12.4, 1.6)
    code_box(s, [
        "IMG_W, IMG_H = 1180.0, 884.0",
        "pixel_to_world(px,py) = (-20.237 + py/IMG_H*40.473,  -26.915 + px/IMG_W*53.83)",
        "world_to_pixel(wx,wy) = ((wy+26.915)/53.83*IMG_W,     (wx+20.237)/40.473*IMG_H)",
    ], top=3.1, size=12.5, height=1.35)
    picture(s, f"{IMG}/track_gt_lane0_demo_render.png", 5.0, 4.5, height=2.5)
    note(s, "▲ 월드좌표 평면에 GT 중심선을 그린 예 (720점, 닫힌 루프)", left=0.6, top=5.4, color=GREY)

    s = add_content(prs, "tools/ — GT 시각화 3종")
    code_box(s, [
        "# (A) 월드좌표 평면 플롯",
        "python3 tools/gt_render.py  ~/track_gt_manual.json -o /tmp/gt_render.png",
        "# (B) 트랙 이미지 위 오버레이",
        "python3 tools/gt_overlay.py ~/track_gt_manual.json -o /tmp/gt_overlay.png",
        "# (C) 실제 GUI 창 스크린샷 (--mask 시 도로/차선 마스킹)",
        "DISPLAY=:1 QT_QPA_PLATFORM=xcb python3 tools/gt_gui_shot.py --mask -o /tmp/gt_mask.png",
    ], top=1.5, size=12, height=2.4)
    picture(s, f"{IMG}/track_gt_lane0_overlay.png", 1.1, 4.15, height=2.35)
    picture(s, f"{IMG}/track_gt_manual_mask_screenshot.png", 7.3, 4.15, height=2.35)
    note(s, "왼쪽: track.png 위 GT 오버레이     |     오른쪽: GUI 도로/차선 마스킹(1차선 빨강·2차선 파랑)",
         top=6.65, color=GREY)

    # ── SECTION 8 ──────────────────────────────────────────────────────────
    add_section(prs, "8. 전체 흐름 & 트러블슈팅")

    s = add_content(prs, "A → Z 전체 흐름 한눈에")
    code_box(s, [
        "[1] teleop_sim.launch.py         시뮬레이터(차+카메라+sender)",
        "[2] teleop_keyboard.py           사람이 키보드로 운전 (w/s/a/d/space/x/q)",
        "[3] manual_collect.py --lane 0   이미지+조작+위치 저장  → manual_demos/",
        "[4] build_wp_from_manual.py      웨이포인트 라벨 생성   → manual_wp_labels.csv",
        "[5] train_head_fast.py           학습(비전고정+Head)    → vla_lora_head_fast.pt",
        "    (또는 train_vla_lora.py       LoRA 정밀학습          → vla_lora_adapter/)",
        "[6] vla_drive.launch.py          학습모델로 자율주행",
        "[7] ros2 topic pub /vla/command  자연어 명령 전달",
        "[8] eval_lap.py / one_lap_test   성능 평가",
    ], top=1.55, size=13.5, height=3.9)
    note(s, "데이터 → 학습 → 추론 → 평가의 한 사이클. 데이터를 늘리고 다시 학습하면 성능이 올라간다.")

    s = add_content(prs, "트러블슈팅 / 함정 노트")
    body(s, [
        ("경로: lora_pipeline 은 ~/ros2_autonomous_vehicle_simulation/lora_pipeline (홈 직속 아님)", 0),
        ("teleop_sim 은 키보드 노드를 켜지 않음 → teleop_keyboard.py 별도 실행", 0),
        ("simulation_pkg 의 data_collection_node 는 소스 없음 → manual_collect.py 사용", 0),
        ("기본 주행은 LoRA가 아니라 '비전고정+spatial head'. LoRA는 VLA_USE_ADAPTER=1로 명시", 0),
        ("'LoRA'가 두 종류: train_lora.py(글자조향) vs train_vla_lora.py(웨이포인트 헤드) — 별개", 0),
        ("VLA 노드는 torch/transformers/peft 필요 (install.sh 에 없음 → 별도 설치)", 0),
        ("teleop_keyboard 의 'c'키 미구현 → 리셋은 x", 0),
        ("gt_annotator.set_show_lane() 인자는 문자열 'both'/'inner'/'outer' (True 넣으면 안 그려짐)", 0),
    ])

    s = add_content(prs, "정리 & 참고 문서")
    body(s, [
        ("VLA 한 줄: 이미지+명령 → (고정 비전 + 작은 학습 헤드) → 웨이포인트 → 조향", 0),
        ("핵심 3원리: ①비전 고정+LoRA  ②웨이포인트 회귀  ③공간정보 보존 헤드", 0),
        ("데이터를 모을수록, 다시 학습할수록 좋아지는 모방학습 사이클", 0),
        ("더 읽을거리", 0),
        ("README.md — Part1(기본) · Part2(VLA 실습 8단계) · 부록A(GT 시각화)", 1),
        ("lora_pipeline/SYSTEM_REPORT.md — Fast/Slow, spatial head 개선기, 추론 상세", 1),
        ("lora_pipeline/DESIGN_pure_vla.md — 순수 단일 VLA 설계 원칙", 1),
        ("tools/README.md — GT 시각화 도구 사용법", 1),
    ])

    add_title_slide(prs, "감사합니다", "Q & A\n자동화연구실 · 성균관대학교")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", default="/home/autolab/Downloads/연구실템플릿 .pptx")
    ap.add_argument("-o", "--out", default="/home/autolab/Downloads/VLA_교육자료.pptx")
    a = ap.parse_args()

    prs = Presentation(a.template)
    _clear_slides(prs)
    build(prs)
    prs.save(a.out)
    print("saved", a.out, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
