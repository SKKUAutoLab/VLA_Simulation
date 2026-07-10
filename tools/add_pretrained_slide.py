#!/usr/bin/env python3
"""사전 학습 산출물 배포 슬라이드를 '사전 준비물' 다음에 끼워 넣는다(제자리 편집)."""
import argparse, copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLUE=RGBColor(0x00,0x70,0xC0); ORANGE=RGBColor(0xED,0x7D,0x31)
BLACK=RGBColor(0x1A,0x1A,0x1A); GREY=RGBColor(0x55,0x55,0x55)
CODEBG=RGBColor(0x1E,0x1E,0x1E); CODEFG=RGBColor(0xE6,0xE6,0xE6); MONO="Consolas"; KO="맑은 고딕"
_B=["❑ ","○ ","· "]; _S=[16,14,13]; _C=[BLACK,RGBColor(0x40,0x40,0x40),RGBColor(0x55,0x55,0x55)]


def _kf(r,sz,bold,color,mono=False):
    r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=color
    r.font.name=MONO if mono else KO
    rPr=r._r.get_or_add_rPr()
    for t in ('a:latin','a:ea','a:cs'):
        e=rPr.find(qn(t))
        if e is None: e=rPr.makeelement(qn(t),{}); rPr.append(e)
        e.set('typeface', MONO if mono else KO)


def title(s, ref, text):
    for ph in ref.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE,PP_PLACEHOLDER.CENTER_TITLE):
            s.shapes._spTree.append(copy.deepcopy(ph._element)); break
    for ph in s.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE,PP_PLACEHOLDER.CENTER_TITLE):
            ph.text_frame.text=text
            for para in ph.text_frame.paragraphs:
                for r in para.runs: r.font.bold=True; r.font.color.rgb=BLUE
            break


def body(s, items, left, top, width, height):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(height))
    tb.text_frame.word_wrap=True; first=True
    for it in items:
        txt,lvl=(it if isinstance(it,tuple) else (it,0)); lvl=min(lvl,2)
        if not txt.strip(): continue
        p=tb.text_frame.paragraphs[0] if first else tb.text_frame.add_paragraph(); first=False
        p.space_after=Pt(5); p.level=lvl
        r=p.add_run(); r.text=("    "*lvl)+_B[lvl]+txt
        _kf(r,_S[lvl],lvl==0,_C[lvl])


def hdr(s,text,top,color=BLUE):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(top),Inches(12.3),Inches(0.4))
    _kf(tb.text_frame.paragraphs[0].add_run(),15,True,color); tb.text_frame.paragraphs[0].runs[0].text=text


def code(s,lines,top,height,left=0.55,width=12.2,size=12.5):
    box=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(left),Inches(top),Inches(width),Inches(height))
    box.fill.solid(); box.fill.fore_color.rgb=CODEBG; box.line.color.rgb=RGBColor(0x40,0x40,0x40)
    box.line.width=Pt(0.75); box.shadow.inherit=False
    tf=box.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.15); tf.vertical_anchor=MSO_ANCHOR.TOP
    first=True
    for ln in lines:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=PP_ALIGN.LEFT; p.space_after=Pt(1)
        r=p.add_run(); r.text=ln or " "
        _kf(r,size,False,RGBColor(0x8A,0xA8,0x8A) if ln.strip().startswith("#") else CODEFG, mono=True)


def note(s,text,top,color=ORANGE,size=12):
    tb=s.shapes.add_textbox(Inches(0.55),Inches(top),Inches(12.2),Inches(0.6)); tb.text_frame.word_wrap=True
    _kf(tb.text_frame.paragraphs[0].add_run(),size,True,color); tb.text_frame.paragraphs[0].runs[0].text=text


def main():
    ap=argparse.ArgumentParser(); ap.add_argument("-f","--file",default="/home/autolab/Downloads/VLA_교육자료.pptx")
    a=ap.parse_args(); prs=Presentation(a.file)
    anchor=None
    for i,s in enumerate(prs.slides):
        if s.shapes and s.shapes[0].has_text_frame and "사전 준비물" in s.shapes[0].text_frame.text:
            anchor=i; break
    if anchor is None: raise SystemExit("'사전 준비물' 슬라이드 못 찾음")
    ref=None
    for s in prs.slides:
        if s.shapes and s.shapes[0].has_text_frame and "핵심 원리" in s.shapes[0].text_frame.text:
            ref=s; break
    ref=ref or prs.slides[anchor]

    s=prs.slides.add_slide(ref.slide_layout)
    title(s, ref, "GPU 부담 줄이기 — 사전 학습 산출물 배포")
    body(s, [
        ("문제: 데이터 수집(66GB)·GPU 학습은 저사양 학생에게 큰 부담", 0),
        ("해결: 학습된 결과물(약 76MB)을 미리 배포 → 수집·학습 단계 건너뛰고 바로 추론", 0),
        ("배포 내용", 0),
        ("vla_lora_head_fast.pt (주행 기본 헤드) · vla_lora_head.pt · vla_lora_adapter/", 1),
        ("배포 제외: 특징 캐시 9.7GB·32GB, 데이터셋 — 용량이 너무 큼", 1),
    ], 0.42, 1.35, 12.4, 2.6)
    hdr(s, "학생 사용법 (데이터·학습 생략)", 3.75)
    code(s, [
        "cd ~/VLA_simulation",
        "bash lora_pipeline/setup_pretrained.sh          # Release에서 산출물 자동 설치",
        "ros2 launch lora_pipeline/vla_drive.launch.py   # 바로 주행",
    ], top=4.2, height=1.35, size=12.5)
    note(s, "⚠️ 배포는 '수집·학습'만 없애줍니다. 추론(주행)은 여전히 GPU 필요 → 없으면 Colab/랩 서버, "
             "또는 데모 영상으로 결과 확인 + 원리 학습.", top=5.75)

    lst=prs.slides._sldIdLst; ids=list(lst); newid=ids[-1]; anchor_id=ids[anchor]
    lst.remove(newid); anchor_id.addnext(newid)
    prs.save(a.file)
    print("inserted pretrained-distribution slide after '사전 준비물' |", len(prs.slides._sldIdLst), "slides")


if __name__=="__main__":
    main()
