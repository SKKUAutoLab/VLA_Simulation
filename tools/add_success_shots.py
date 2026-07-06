#!/usr/bin/env python3
"""
add_success_shots.py — '설치·실행 성공 화면' 슬라이드(실제 스크린샷)를
'설치 & 빌드' 다음에 끼워 넣는다. 재생성 없이 제자리 편집(다른 슬라이드 보존).
"""
import argparse, copy, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

BLUE=RGBColor(0x00,0x70,0xC0); ORANGE=RGBColor(0xED,0x7D,0x31); GREY=RGBColor(0x55,0x55,0x55)
KO="맑은 고딕"
VERIFY="/home/autolab/shot_verify.png"
GAZEBO="/home/autolab/shot_gazebo.png"


def _kf(r, sz, bold, color):
    r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=color; r.font.name=KO
    rPr=r._r.get_or_add_rPr()
    for t in ('a:latin','a:ea','a:cs'):
        e=rPr.find(qn(t))
        if e is None: e=rPr.makeelement(qn(t),{}); rPr.append(e)
        e.set('typeface',KO)


def _title(s, ref, title):
    for ph in ref.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            s.shapes._spTree.append(copy.deepcopy(ph._element)); break
    for ph in s.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            ph.text_frame.text=title
            for para in ph.text_frame.paragraphs:
                for r in para.runs: r.font.bold=True; r.font.color.rgb=BLUE
            break


def cap(s, text, left, top, width, color=GREY, size=12, bold=True):
    tb=s.shapes.add_textbox(Inches(left),Inches(top),Inches(width),Inches(0.35))
    tb.text_frame.word_wrap=True
    p=tb.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    _kf(p.add_run(),size,bold,color); p.runs[0].text=text


def pic_framed(s, path, left, top, width):
    from PIL import Image
    w,h=Image.open(path).size
    height=width*h/w
    fr=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(left-0.03),Inches(top-0.03),
                          Inches(width+0.06),Inches(height+0.06))
    fr.fill.background(); fr.line.color.rgb=RGBColor(0x88,0x88,0x88); fr.line.width=Pt(1)
    fr.shadow.inherit=False
    s.shapes.add_picture(path,Inches(left),Inches(top),width=Inches(width))
    return height


def main():
    ap=argparse.ArgumentParser()
    ap.add_argument("-f","--file",default="/home/autolab/Downloads/VLA_교육자료.pptx")
    a=ap.parse_args()
    for p in (VERIFY,GAZEBO):
        if not os.path.exists(p): raise SystemExit("이미지 없음: "+p)
    prs=Presentation(a.file)

    # '설치 & 빌드' 슬라이드 인덱스
    anchor=None
    for i,s in enumerate(prs.slides):
        if s.shapes and s.shapes[0].has_text_frame and "설치 & 빌드" in s.shapes[0].text_frame.text:
            anchor=i; break
    if anchor is None: raise SystemExit("'설치 & 빌드' 슬라이드 못 찾음")
    ref=prs.slides[anchor]  # 제목 틀 복제용

    s=prs.slides.add_slide(ref.slide_layout)
    _title(s, ref, "설치·실행 성공 화면 — 이렇게 나오면 정상")

    # 좌: Gazebo 시뮬레이터 실행 성공
    cap(s,"① 시뮬레이터 실행 성공  (ros2 launch simulation_pkg driving_sim.launch.py)",0.35,1.35,7.4,BLUE,12.5)
    pic_framed(s, GAZEBO, 0.5, 1.8, 7.1)
    cap(s,"트랙 위에 차량이 스폰되고 Gazebo 3D 화면이 뜨면 성공",0.5,6.55,7.1,GREY,11)

    # 우: 환경 검증 터미널
    cap(s,"② 환경 검증 성공  (GPU · CUDA · ROS2)",8.0,1.35,5.0,BLUE,12.5)
    h=pic_framed(s, VERIFY, 8.0, 1.8, 4.95)
    cap(s,"CUDA: True 와 GPU 이름이 나오면 준비 완료",8.0,1.8+h+0.1,4.95,GREY,11)
    # 하단 안내
    tb=s.shapes.add_textbox(Inches(8.0),Inches(4.4),Inches(5.0),Inches(2.0))
    tb.text_frame.word_wrap=True
    for i,line in enumerate([
        "확인 포인트",
        "· nvidia-smi 에 GPU 표가 뜬다",
        "· torch.cuda.is_available() == True",
        "· ROS_DISTRO = humble",
        "· Gazebo 창에 트랙과 차량이 보인다",
    ]):
        p=tb.text_frame.paragraphs[0] if i==0 else tb.text_frame.add_paragraph()
        _kf(p.add_run(), 13 if i==0 else 12, i==0, RGBColor(0x1A,0x1A,0x1A) if i==0 else GREY)
        p.runs[0].text=line; p.space_after=Pt(3)

    # 삽입: 설치&빌드 바로 뒤
    lst=prs.slides._sldIdLst; ids=list(lst)
    newid=ids[-1]; anchor_id=ids[anchor]
    lst.remove(newid); anchor_id.addnext(newid)
    prs.save(a.file)
    print("inserted success-shots slide after '설치 & 빌드' |", len(prs.slides._sldIdLst), "slides")


if __name__=="__main__":
    main()
